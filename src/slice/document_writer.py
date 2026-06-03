"""Document writing utilities for various file formats."""

from pathlib import Path
from typing import Dict, Any, List, Union


def write_document(file_path: str, operations: Union[Dict, List[Dict]]) -> Dict[str, Any]:
    """
    Write/modify documents in various formats.

    Supported formats:
    - Word (.docx) - append paragraphs, replace text, modify tables
    - Excel (.xlsx) - set cells, append rows, modify columns
    - PowerPoint (.pptx) - add slides, modify slide text
    - CSV (.csv) - append rows, modify cells
    - Text files (.txt, .md, etc.) - replace content, append text
    - PDF (.pdf) - create pages, add text/paragraphs, merge PDFs

    Args:
        file_path: Path to the document to write/modify
        operations: Single operation dict or list of operations to perform.
                   Each operation has 'type' and operation-specific parameters.

    Returns:
        Dict with 'success', 'message', 'error', 'operations_applied' keys

    Operation Examples:

    Word (.docx):
        {"type": "append_paragraph", "text": "New paragraph text"}
        {"type": "replace_text", "find": "old text", "replace": "new text"}
        {"type": "insert_after", "search": "Section Header", "text": "New content"}

    Excel (.xlsx):
        {"type": "set_cell", "sheet": "Sheet1", "row": 5, "col": 3, "value": "Data"}
        {"type": "append_row", "sheet": "Sheet1", "values": ["A", "B", "C"]}
        {"type": "set_column", "sheet": "Sheet1", "col": "M", "start_row": 3, "values": ["X", "Y", "Z"]}

    PowerPoint (.pptx):
        {"type": "add_slide", "title": "Slide Title", "content": "Slide content text"}

    PDF (.pdf):
        {"type": "add_page", "title": "Page Title", "content": "Page content text"}
        {"type": "add_text", "text": "Text content", "x": 100, "y": 500, "font_size": 12}
        {"type": "add_paragraph", "text": "Paragraph text", "font_size": 12}

    CSV (.csv):
        {"type": "append_row", "values": ["col1", "col2", "col3"]}
        {"type": "set_cell", "row": 2, "col": 1, "value": "New Value"}

    Text files:
        {"type": "replace_content", "text": "Entirely new content"}
        {"type": "append_text", "text": "\\nAppended text"}
    """
    path = Path(file_path)

    # Normalize operations to list
    if isinstance(operations, dict):
        operations = [operations]

    # Check if file exists for modification operations
    # (Some operations like creating new files don't require existing file)
    file_exists = path.exists()

    # Determine file type
    suffix = path.suffix.lower()

    try:
        if suffix == '.pdf':
            result = _write_pdf(path, operations, file_exists)
        elif suffix == '.docx':
            result = _write_docx(path, operations, file_exists)
        elif suffix == '.xlsx':
            result = _write_excel(path, operations, file_exists)
        elif suffix == '.pptx':
            result = _write_pptx(path, operations, file_exists)
        elif suffix == '.csv':
            result = _write_csv(path, operations, file_exists)
        else:
            # Treat as plain text
            result = _write_text(path, operations, file_exists)

        return result

    except Exception as e:
        return {
            "success": False,
            "message": "",
            "error": f"Failed to write {file_path}: {str(e)}",
            "operations_applied": 0
        }


def _write_docx(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/modify Word documents."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required to write Word documents. "
            "Please reinstall slice-agent: pip install -e ."
        )

    # Load existing or create new document
    if file_exists:
        doc = Document(path)
    else:
        doc = Document()

    operations_applied = 0
    messages = []

    for op in operations:
        op_type = op.get("type", "")

        if op_type == "append_paragraph":
            text = op.get("text", "")
            doc.add_paragraph(text)
            operations_applied += 1
            messages.append(f"Appended paragraph with {len(text)} characters")

        elif op_type == "replace_text":
            find = op.get("find", "")
            replace = op.get("replace", "")
            count = 0
            for paragraph in doc.paragraphs:
                if find in paragraph.text:
                    paragraph.text = paragraph.text.replace(find, replace)
                    count += 1
            # Also check tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if find in cell.text:
                            cell.text = cell.text.replace(find, replace)
                            count += 1
            operations_applied += 1
            messages.append(f"Replaced '{find}' with '{replace}' in {count} locations")

        elif op_type == "insert_after":
            search = op.get("search", "")
            text = op.get("text", "")
            found = False
            for i, paragraph in enumerate(doc.paragraphs):
                if search in paragraph.text:
                    # Insert new paragraph after this one
                    # This is tricky in python-docx, we'll append and note the limitation
                    doc.add_paragraph(text)
                    found = True
                    operations_applied += 1
                    messages.append(f"Inserted text after '{search}' (note: added at end due to python-docx limitations)")
                    break
            if not found:
                messages.append(f"Warning: Could not find '{search}' to insert after")

        else:
            messages.append(f"Unknown operation type: {op_type}")

    # Save the document
    doc.save(path)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }


def _write_excel(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/modify Excel spreadsheets."""
    try:
        from openpyxl import load_workbook, Workbook
        from openpyxl.utils import get_column_letter, column_index_from_string
    except ImportError:
        raise ImportError(
            "openpyxl is required to write Excel files. "
            "Please reinstall slice-agent: pip install -e ."
        )

    # Load existing or create new workbook
    if file_exists:
        # data_only=True evaluates formulas to their values
        workbook = load_workbook(path, data_only=True)
    else:
        workbook = Workbook()
        # Remove default sheet if it exists and we'll create custom sheets
        # Default sheet name varies by version, so check all existing sheets
        if workbook.sheetnames and len(workbook.sheetnames) == 1:
            # Only one default sheet exists, will be replaced when we create_sheet
            pass

    operations_applied = 0
    messages = []

    for op in operations:
        op_type = op.get("type", "")
        sheet_name = op.get("sheet", workbook.sheetnames[0] if workbook.sheetnames else "Sheet1")

        # Get or create sheet
        if sheet_name not in workbook.sheetnames:
            sheet = workbook.create_sheet(sheet_name)
        else:
            sheet = workbook[sheet_name]

        if op_type == "set_cell":
            row = op.get("row", 1)
            col = op.get("col", 1)
            value = op.get("value", "")

            # Handle column as letter or number
            if isinstance(col, str):
                col_idx = column_index_from_string(col)
            else:
                col_idx = col

            sheet.cell(row=row, column=col_idx, value=value)
            operations_applied += 1
            messages.append(f"Set cell {get_column_letter(col_idx)}{row} = '{value}'")

        elif op_type == "append_row":
            values = op.get("values", [])
            sheet.append(values)
            operations_applied += 1
            messages.append(f"Appended row with {len(values)} values")

        elif op_type == "set_column":
            col = op.get("col", 1)
            start_row = op.get("start_row", 1)
            values = op.get("values", [])

            # Handle column as letter or number
            if isinstance(col, str):
                col_idx = column_index_from_string(col)
            else:
                col_idx = col

            for i, value in enumerate(values):
                sheet.cell(row=start_row + i, column=col_idx, value=value)

            operations_applied += 1
            messages.append(f"Set column {get_column_letter(col_idx)} starting at row {start_row} with {len(values)} values")

        else:
            messages.append(f"Unknown operation type: {op_type}")

    # Save the workbook
    workbook.save(path)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }


def _write_pptx(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/modify PowerPoint presentations."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required to write PowerPoint files. "
            "Please reinstall slice-agent: pip install -e ."
        )

    # Load existing or create new presentation
    if file_exists:
        prs = Presentation(path)
    else:
        prs = Presentation()

    operations_applied = 0
    messages = []

    for op in operations:
        op_type = op.get("type", "")

        if op_type == "add_slide":
            title = op.get("title", "")
            content = op.get("content", "")

            # Use blank slide layout (index 6 is typically blank)
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title

            # Add content to first text placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Body placeholder
                    shape.text = content
                    break

            operations_applied += 1
            messages.append(f"Added slide: '{title}'")

        else:
            messages.append(f"Unknown operation type: {op_type}")

    # Save the presentation
    prs.save(path)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }


def _write_pdf(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/create PDF documents."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_LEFT
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        raise ImportError(
            "reportlab and pypdf are required to write PDF files. "
            "Please reinstall slice: pip install -e ."
        )

    operations_applied = 0
    messages = []

    # For creating new PDFs or adding content
    if not file_exists or any(op.get("type") in ["add_page", "add_paragraph", "add_text"] for op in operations):
        # Create a new PDF with content
        doc = SimpleDocTemplate(str(path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "add_page":
                title = op.get("title", "")
                content = op.get("content", "")

                # Add page break if not first page
                if operations_applied > 0:
                    story.append(PageBreak())

                # Add title
                if title:
                    title_style = ParagraphStyle(
                        'CustomTitle',
                        parent=styles['Heading1'],
                        fontSize=24,
                        spaceAfter=30
                    )
                    story.append(Paragraph(title, title_style))

                # Add content
                if content:
                    story.append(Paragraph(content, styles['BodyText']))
                    story.append(Spacer(1, 0.2*inch))

                operations_applied += 1
                messages.append(f"Added page: '{title}'")

            elif op_type == "add_paragraph":
                text = op.get("text", "")
                font_size = op.get("font_size", 12)

                para_style = ParagraphStyle(
                    'CustomPara',
                    parent=styles['BodyText'],
                    fontSize=font_size,
                    spaceAfter=12
                )
                story.append(Paragraph(text, para_style))
                operations_applied += 1
                messages.append(f"Added paragraph with {len(text)} characters")

            elif op_type == "add_text":
                text = op.get("text", "")
                font_size = op.get("font_size", 12)

                # For simple text, treat as paragraph
                para_style = ParagraphStyle(
                    'CustomText',
                    parent=styles['Normal'],
                    fontSize=font_size,
                    alignment=TA_LEFT
                )
                story.append(Paragraph(text, para_style))
                story.append(Spacer(1, 0.1*inch))
                operations_applied += 1
                messages.append(f"Added text: '{text[:50]}...' " if len(text) > 50 else f"Added text: '{text}'")

            else:
                messages.append(f"Unknown operation type: {op_type}")

        # Build the PDF
        doc.build(story)

    # For merging PDFs
    elif file_exists and any(op.get("type") == "merge_pdf" for op in operations):
        writer = PdfWriter()

        # Read existing PDF
        reader = PdfReader(str(path))
        for page in reader.pages:
            writer.add_page(page)

        for op in operations:
            op_type = op.get("type", "")

            if op_type == "merge_pdf":
                source_path = op.get("source", "")
                if source_path and Path(source_path).exists():
                    source_reader = PdfReader(source_path)
                    for page in source_reader.pages:
                        writer.add_page(page)
                    operations_applied += 1
                    messages.append(f"Merged PDF from '{source_path}'")
                else:
                    messages.append(f"Warning: Source PDF not found: '{source_path}'")

        # Write the merged PDF
        with open(path, 'wb') as output_file:
            writer.write(output_file)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }


def _write_csv(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/modify CSV files."""
    import csv

    # Read existing content if file exists
    rows = []
    if file_exists:
        # Try UTF-8 first, then latin-1 for encoding compatibility
        last_error = None
        for encoding in ['utf-8', 'latin-1']:
            try:
                with open(path, 'r', newline='', encoding=encoding) as f:
                    reader = csv.reader(f)
                    rows = list(reader)
                break  # Success, exit encoding loop
            except UnicodeDecodeError as e:
                last_error = e
                continue
            except Exception:
                # Other errors (permissions, etc.) should be raised
                raise
        else:
            # If we exhausted all encodings, raise the last error
            if last_error:
                raise last_error

    operations_applied = 0
    messages = []

    for op in operations:
        op_type = op.get("type", "")

        if op_type == "append_row":
            values = op.get("values", [])
            rows.append(values)
            operations_applied += 1
            messages.append(f"Appended row with {len(values)} values")

        elif op_type == "set_cell":
            row_idx = op.get("row", 1) - 1  # Convert to 0-indexed
            col_idx = op.get("col", 1) - 1  # Convert to 0-indexed
            value = op.get("value", "")

            # Ensure row exists
            while len(rows) <= row_idx:
                rows.append([])

            # Ensure column exists in row
            while len(rows[row_idx]) <= col_idx:
                rows[row_idx].append("")

            rows[row_idx][col_idx] = value
            operations_applied += 1
            messages.append(f"Set cell (row {row_idx + 1}, col {col_idx + 1}) = '{value}'")

        else:
            messages.append(f"Unknown operation type: {op_type}")

    # Write back to file
    with open(path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerows(rows)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }


def _write_text(path: Path, operations: List[Dict], file_exists: bool) -> Dict[str, Any]:
    """Write/modify plain text files."""
    # Read existing content if file exists
    content = ""
    if file_exists:
        # Try UTF-8 first, then latin-1, then binary with error replacement
        try:
            with open(path, 'r', encoding='utf-8') as f:
                content = f.read()
        except UnicodeDecodeError:
            try:
                with open(path, 'r', encoding='latin-1') as f:
                    content = f.read()
            except Exception:
                # Last resort: binary read with error replacement
                with open(path, 'rb') as f:
                    content = f.read().decode('utf-8', errors='replace')

    operations_applied = 0
    messages = []

    for op in operations:
        op_type = op.get("type", "")

        if op_type == "replace_content":
            text = op.get("text", "")
            content = text
            operations_applied += 1
            messages.append(f"Replaced entire content with {len(text)} characters")

        elif op_type == "append_text":
            text = op.get("text", "")
            content += text
            operations_applied += 1
            messages.append(f"Appended {len(text)} characters")

        elif op_type == "replace_text":
            find = op.get("find", "")
            replace = op.get("replace", "")
            count = content.count(find)
            content = content.replace(find, replace)
            operations_applied += 1
            messages.append(f"Replaced '{find}' with '{replace}' ({count} occurrences)")

        else:
            messages.append(f"Unknown operation type: {op_type}")

    # Write back to file
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)

    return {
        "success": True,
        "message": "; ".join(messages),
        "error": "",
        "operations_applied": operations_applied
    }
